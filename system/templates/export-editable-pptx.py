#!/usr/bin/env python3
"""排版小幫手輸出器：把輪播頁 HTML 轉成可編輯 PPTX（真文字框，非圖片）。

原理：用 Playwright 讀每頁 HTML，走 DOM 找「文字單元」（葉節點，或子節點
全是行內標籤 span/br 等的節點——碰到區塊邊界就停），記下每個單元的座標、
對齊、逐 run 樣式（顏色／背景／粗細／字級／字型），以及是單行還是多行；
接著把這些文字單元隱藏、對 .page 截圖，拿到「無字背景」；最後用 python-pptx
組 PPTX：背景若接近純色就用原生色塊填充，有真實圖形內容（如岔路圖線條）
才嵌圖；文字若帶有實色背景（藥丸／卡片）就轉成圓角矩形圖案，其餘轉成透明
文字框。全程只在記憶體處理背景截圖，不寫暫存檔到來源資料夾。

螢光筆：不用 PPTX 的文字高亮屬性（<a:highlight>）——這個屬性 PowerPoint 認得，
但 Keynote 的 pptx 匯入不支援，開起來會整個消失。改用「畫一塊黃色色塊放在文字
底下」的圖形做法，Keynote／PowerPoint／Google Slides 都能正確顯示。

用法：
  <裝有 playwright + python-pptx + Pillow 的 python3> export-editable-pptx.py \\
    P1.html P2.html ... [--out 檔名.pptx]

  - 不給 --out 時，預設輸出到第一個輸入檔同層目錄的「可編輯版.pptx」
  - 字型需求：本機要裝好 jf-jinxuan 各字重的 OpenType 家族（見下方
    JINXUAN_WEIGHTS），沒裝的話 PowerPoint/Keynote 開啟時會退回系統預設字型
"""
import os, sys, re, io
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from PIL import Image

# ---------------------------------------------------------------------------
# CLI 參數（手動解析，不用 argparse，跟 export.py 同風格）
# ---------------------------------------------------------------------------
args = list(sys.argv[1:])
out_path = None
if "--out" in args:
    i = args.index("--out"); out_path = args[i + 1]; del args[i:i + 2]
files = args
if not files:
    sys.exit("用法：export-editable-pptx.py <頁面.html ...> [--out 檔名.pptx]")

if out_path is None:
    out_path = os.path.join(os.path.dirname(os.path.abspath(files[0])), "可編輯版.pptx")

# ---------------------------------------------------------------------------
# 字型：實際安裝字型的內部家族名（每個字重是獨立家族，不是同一家族的粗細變化）
# ---------------------------------------------------------------------------
JINXUAN_WEIGHTS = {
    900: "jf-jinxuan ExtraBold",
    700: "jf-jinxuan Bold",
    500: "jf-jinxuan Medium",
    400: "jf-jinxuan Book",
    200: "jf-jinxuan UltraLight",
}

def jinxuan_font_name(weight):
    nearest = min(JINXUAN_WEIGHTS.keys(), key=lambda w: abs(w - weight))
    return JINXUAN_WEIGHTS[nearest]

# ---------------------------------------------------------------------------
# 「.body」段落文字的特殊排版覆寫（2026-07-28 使用者在 Keynote 手動檢視後定案）
# 只套用在 class list 包含 "body" 的文字單元（敘事段落），其餘單元類型
# （大標／說明字／頁尾／藥丸／引言……）完全不受影響。
# ---------------------------------------------------------------------------
BODY_FONT_PT = 28          # 蓋掉 px*0.75 推算值（2026-07-28 使用者定案：主要內文 28pt）
BODY_FONT_PX = BODY_FONT_PT / 0.75   # 換算回瀏覽器量測用的 px（96dpi 假設，全檔案一致）
BODY_LINE_SPACING = 1.3    # 段落行距
BODY_LETTER_SPACING = 0.13 # 字距，字級的比例（Keynote 顯示 13%；2026-07-28 使用者定案）
BODY_BOX_W_PX = 850        # 窄化置中框寬
BODY_BOX_X_PX = 115        # (1080-850)/2，左右對稱留白

HIGHLIGHT_YELLOW = (247, 233, 78)

# 插畫位／照片位（class 含 illust-slot 或 photo-slot）的呈現規則
# （2026-07-29 使用者定案：直角、灰色虛線邊框、無陰影、25pt、文字上下置中、單一物件。
#  直角的原因：Keynote 匯入器遇到「圓角＋可見底色或框線」必拆成兩層群組——實驗證實
#  連 Keynote 自家輸出的 pptx 重新匯入也一樣，無 XML 寫法可繞過；插畫位終會被生成圖
#  蓋掉，捨圓角換單一物件。膠囊與句級螢光框則反向取捨：保留圓角、接受群組）
SLOT_FONT_PT = 25
SLOT_BORDER_GRAY = (138, 138, 133)   # 同插畫位說明文字的灰

# 句級螢光筆（.hl-line 黃色圓角框）的呈現規則
# （2026-07-28 使用者定案：文字直接打在圓角框內、上下置中、行距 1.3、字距 13%，單一物件）
HLLINE_LINE_SPACING = 1.3
HLLINE_LETTER_SPACING = 0.13
HLLINE_RADIUS_PX = 8                 # 同 CSS .hl-line 的 border-radius
HLLINE_PAD_V_PT = 8                  # 黃框上下各加的內距（2026-07-29 使用者定案）

# 滿版頁（內容從頁頂邊距就開始的頁，如故事頁/收尾頁；封面除外）第一個物件的
# 起始 y 統一下移到 80pt（2026-07-29 使用者定案；HTML 頁邊距 72px=54pt 太貼頂）
FULLBLEED_FIRST_Y_PT = 80

def dominant_color_deviation(img):
    # 用原尺寸判斷（不能縮圖，縮圖會把細線圖形抗鋸齒掉，誤判成純色）
    img = img.convert("RGB")
    colors = img.getcolors(maxcolors=img.width * img.height)
    colors.sort(reverse=True)
    total = sum(c for c, _ in colors)
    top_count, top_color = colors[0]
    return top_color, total - top_count

PX_W, PX_H = 1080, 1350
EMU_PER_PX = 914400 / 96

def px(v):
    return Emu(int(round(v * EMU_PER_PX)))

def rgb_str_to_hex(s):
    m = re.match(r"rgba?\(([\d.]+),\s*([\d.]+),\s*([\d.]+)", s)
    if not m:
        return None
    r, g, b = (int(round(float(x))) for x in m.groups())
    return RGBColor(r, g, b)

def is_transparent(s):
    return s.startswith("rgba(0, 0, 0, 0)") or s == "transparent"

def is_highlight(s):
    c = rgb_str_to_hex(s)
    if c is None:
        return False
    return (c[0], c[1], c[2]) == HIGHLIGHT_YELLOW

def set_cjk_font(run, name):
    # PowerPoint 用「東亞字型」欄位（a:ea）渲染中文，只設 a:latin 中文不會生效
    rPr = run._r.get_or_add_rPr()
    latin = rPr.find(qn('a:latin'))
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = etree_element(qn('a:ea'))
        if latin is not None:
            latin.addnext(ea)
        else:
            rPr.append(ea)
    ea.set('typeface', name)

from lxml import etree as _etree
def etree_element(tag):
    return _etree.Element(tag)

def align_of(css_align):
    return {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT,
            "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}.get(css_align, PP_ALIGN.LEFT)

def set_round_rect(box, adj_val):
    # 把 textbox 的幾何改成圓角矩形。文字直接放在帶底色的 textbox 裡（單一物件），
    # 不用「autoshape 圖形＋文字」——Keynote 匯入 pptx 時會把後者拆成兩層。
    # adj_val：OOXML roundRect 的圓角參數，圓角半徑 = min(寬,高) × adj_val/100000。
    # 注意 spPr 子元素順序（xfrm → prstGeom → fill → ln）,幾何必須插在 xfrm 之後、
    # 設定 fill/line 之前，順序錯了檔案有被判定損毀的風險。
    spPr = box._element.spPr
    for g in spPr.findall(qn('a:prstGeom')):
        spPr.remove(g)
    geom = _etree.Element(qn('a:prstGeom'))
    geom.set('prst', 'roundRect')
    av = _etree.SubElement(geom, qn('a:avLst'))
    gd = _etree.SubElement(av, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', 'val %d' % int(adj_val))
    xfrm = spPr.find(qn('a:xfrm'))
    if xfrm is not None:
        xfrm.addnext(geom)
    else:
        spPr.insert(0, geom)

def add_highlight_shape(slide, x_px, y_px, w_px, h_px):
    # 螢光筆改用圖形色塊（不用 <a:highlight> 文字屬性——Keynote 的 pptx 匯入不支援該屬性）
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  px(x_px), px(y_px), px(w_px), px(h_px))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(*HIGHLIGHT_YELLOW)
    shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.18
    except Exception:
        pass
    return shp

# ---------------------------------------------------------------------------
# 瀏覽器端擷取邏輯：找文字單元、量座標樣式、量螢光筆色塊矩形
# ---------------------------------------------------------------------------
EXTRACT_JS = """
() => {
  const INLINE_TAGS = new Set(['SPAN','BR','I','B','EM','STRONG','A']);
  function isInlineOnly(el) {
    for (const child of Array.from(el.children)) {
      if (!INLINE_TAGS.has(child.tagName)) return false;
      if (!isInlineOnly(child)) return false;
    }
    return true;
  }
  function isHighlightBg(bgStr) {
    const m = bgStr.match(/rgba?\\((\\d+),\\s*(\\d+),\\s*(\\d+)/);
    if (!m) return false;
    return parseInt(m[1]) === 247 && parseInt(m[2]) === 233 && parseInt(m[3]) === 78;
  }
  function styleOf(node) {
    const cs = getComputedStyle(node);
    return {
      color: cs.color,
      bg: cs.backgroundColor,
      weight: parseInt(cs.fontWeight) || 400,
      fontSize: parseFloat(cs.fontSize),
      fontFamily: cs.fontFamily.split(',')[0].replace(/["']/g, '').trim(),
    };
  }
  // 依 display 決定螢光筆色塊怎麼量：inline 元素會跨行斷成好幾塊（box-decoration-break:
  // clone 的效果，每行各自一個圓角色塊）；inline-block/block 是單一整塊（如 .hl-line）
  function highlightRectsOf(el, originX, originY) {
    const disp = getComputedStyle(el).display;
    const out = [];
    if (disp === 'inline') {
      const range = document.createRange();
      range.selectNodeContents(el);
      const raw = Array.from(range.getClientRects()).filter(r => r.width > 0 && r.height > 0);
      const groups = [];
      for (const rc of raw) {
        let g = groups.find(g => Math.abs(g.top - rc.top) < 3);
        if (!g) { g = { top: rc.top, bottom: rc.bottom, left: rc.left, right: rc.right }; groups.push(g); }
        else {
          g.top = Math.min(g.top, rc.top); g.bottom = Math.max(g.bottom, rc.bottom);
          g.left = Math.min(g.left, rc.left); g.right = Math.max(g.right, rc.right);
        }
      }
      for (const g of groups) {
        out.push({ x: g.left - originX, y: g.top - originY, w: g.right - g.left, h: g.bottom - g.top });
      }
    } else {
      const r = el.getBoundingClientRect();
      out.push({ x: r.left - originX, y: r.top - originY, w: r.width, h: r.height });
    }
    return out;
  }
  function getParagraphs(el, pageRect) {
    const paragraphs = [[]];
    function walk(node, styleEl, hlRects) {
      for (const child of Array.from(node.childNodes)) {
        if (child.nodeType === Node.TEXT_NODE) {
          const t = child.textContent;
          if (t && t.trim() !== '') {
            paragraphs[paragraphs.length - 1].push({ text: t, ...styleOf(styleEl), hlRects: hlRects || null });
          }
        } else if (child.nodeType === Node.ELEMENT_NODE) {
          if (child.tagName === 'BR') {
            paragraphs.push([]);
          } else {
            const cs = getComputedStyle(child);
            let childHl = hlRects;
            if (isHighlightBg(cs.backgroundColor)) {
              childHl = highlightRectsOf(child, pageRect.left, pageRect.top);
            }
            walk(child, child, childHl);
          }
        }
      }
    }
    // 單元根節點本身就是螢光背景時（例如單元就是 .hl-line span），也要量到色塊矩形
    const rootHl = isHighlightBg(getComputedStyle(el).backgroundColor)
      ? highlightRectsOf(el, pageRect.left, pageRect.top) : null;
    walk(el, el, rootHl);
    return paragraphs.filter((p, i) => p.length > 0 || i === paragraphs.length - 1);
  }
  const page = document.querySelector('.page');
  const pageRect = page.getBoundingClientRect();
  const units = [];
  function walk(el) {
    if (el.children.length === 0) {
      if (el.textContent && el.textContent.trim() !== '') units.push(el);
      return;
    }
    if (isInlineOnly(el) && el.textContent && el.textContent.trim() !== '') {
      units.push(el);
      return;
    }
    for (const child of Array.from(el.children)) walk(child);
  }
  walk(page);
  function countLines(el) {
    const range = document.createRange();
    range.selectNodeContents(el);
    const rects = Array.from(range.getClientRects());
    const tops = [];
    for (const rc of rects) {
      if (rc.width === 0 || rc.height === 0) continue;
      const t = Math.round(rc.top);
      if (!tops.some(x => Math.abs(x - t) < 3)) tops.push(t);
    }
    return Math.max(tops.length, 1);
  }
  const result = units.map(el => {
    const r = el.getBoundingClientRect();
    const cs = getComputedStyle(el);
    return {
      x: r.left - pageRect.left,
      y: r.top - pageRect.top,
      w: r.width,
      h: r.height,
      textAlign: cs.textAlign,
      lines: countLines(el),
      className: el.className || '',
      hasHlLine: el.classList.contains('hl-line') || !!el.querySelector('.hl-line'),
      outerHTML: el.outerHTML,
      paragraphs: getParagraphs(el, pageRect),
    };
  });
  units.forEach(el => { el.style.visibility = 'hidden'; });
  return result;
}
"""

# 針對 .body 覆寫框（窄化＋放大字級＋加行距）重新量測：字級/框寬一變，原本抓到的
# 螢光筆色塊座標、以及原始高度就對不上新版面了，所以用同一份 outerHTML 在一個
# 隱藏容器裡套上新的 width/font-size/line-height 重新排一次版、重新量。
REMEASURE_JS = """
({ html, boxWpx, fontPx, lineHeight, letterSpacingEm }) => {
  function isHighlightBg(bgStr) {
    const m = bgStr.match(/rgba?\\((\\d+),\\s*(\\d+),\\s*(\\d+)/);
    if (!m) return false;
    return parseInt(m[1]) === 247 && parseInt(m[2]) === 233 && parseInt(m[3]) === 78;
  }
  function highlightRectsOf(el, originX, originY) {
    const disp = getComputedStyle(el).display;
    const out = [];
    if (disp === 'inline') {
      const range = document.createRange();
      range.selectNodeContents(el);
      const raw = Array.from(range.getClientRects()).filter(r => r.width > 0 && r.height > 0);
      const groups = [];
      for (const rc of raw) {
        let g = groups.find(g => Math.abs(g.top - rc.top) < 3);
        if (!g) { g = { top: rc.top, bottom: rc.bottom, left: rc.left, right: rc.right }; groups.push(g); }
        else {
          g.top = Math.min(g.top, rc.top); g.bottom = Math.max(g.bottom, rc.bottom);
          g.left = Math.min(g.left, rc.left); g.right = Math.max(g.right, rc.right);
        }
      }
      for (const g of groups) out.push({ x: g.left - originX, y: g.top - originY, w: g.right - g.left, h: g.bottom - g.top });
    } else {
      const r = el.getBoundingClientRect();
      out.push({ x: r.left - originX, y: r.top - originY, w: r.width, h: r.height });
    }
    return out;
  }
  const container = document.createElement('div');
  container.style.cssText = `position:fixed; left:-9999px; top:0; width:${boxWpx}px;`;
  container.innerHTML = html;
  document.body.appendChild(container);
  const inner = container.firstElementChild;
  inner.style.width = boxWpx + 'px';
  inner.style.fontSize = fontPx + 'px';
  inner.style.lineHeight = String(lineHeight);
  inner.style.letterSpacing = letterSpacingEm + 'em';
  const containerRect = container.getBoundingClientRect();
  const hlEls = Array.from(inner.querySelectorAll('*')).filter(
    e => isHighlightBg(getComputedStyle(e).backgroundColor)
  );
  const hlRects = [];
  for (const e of hlEls) hlRects.push(...highlightRectsOf(e, containerRect.left, containerRect.top));
  const height = inner.getBoundingClientRect().height;
  document.body.removeChild(container);
  return { height, hlRects };
}
"""

# SVG 圖表（循環箭頭、岔路線…）獨立擷取（2026-07-29 使用者定案：不再烙進背景圖，
# 改成透明底 PNG 獨立物件，交付後可在 Keynote 單獨移動）。
# 步驟：把 svg 的祖先鏈與頁面背景暫設透明 → 逐個 svg 以 omit_background 截透明圖
# → 還原頁面背景 → 把 svg 隱藏（之後的背景截圖就不含圖表，多半能回到純色底）
SVG_PREP_JS = """
() => {
  const page = document.querySelector('.page');
  const pr = page.getBoundingClientRect();
  const svgs = Array.from(page.querySelectorAll('svg'));
  const out = [];
  svgs.forEach((s, i) => {
    s.setAttribute('data-export-svg', String(i));
    let n = s.parentElement;
    while (n && n !== page) { n.style.background = 'transparent'; n = n.parentElement; }
    const r = s.getBoundingClientRect();
    out.push({ x: r.left - pr.left, y: r.top - pr.top, w: r.width, h: r.height });
  });
  if (svgs.length) {
    page.style.background = 'transparent';
    document.body.style.background = 'transparent';
    document.documentElement.style.background = 'transparent';
  }
  return out;
}
"""
SVG_RESTORE_HIDE_JS = """
() => {
  const page = document.querySelector('.page');
  page.style.background = '';
  document.body.style.background = '';
  document.documentElement.style.background = '';
  page.querySelectorAll('svg').forEach(s => { s.style.visibility = 'hidden'; });
}
"""

# 句級螢光筆（.hl-line）重量測：套上新字距/行距後，黃框的自然寬高會變
# （inline-block 寬度跟著內容長），量出新尺寸供置中擺放
HLLINE_REMEASURE_JS = """
({ html, letterSpacingEm, lineHeight }) => {
  const container = document.createElement('div');
  container.style.cssText = 'position:fixed; left:-9999px; top:0;';
  container.innerHTML = html;
  document.body.appendChild(container);
  const root = container.firstElementChild;
  const target = root.classList.contains('hl-line') ? root : root.querySelector('.hl-line');
  target.style.letterSpacing = letterSpacingEm + 'em';
  target.style.lineHeight = String(lineHeight);
  const r = target.getBoundingClientRect();
  const res = { w: r.width, h: r.height };
  document.body.removeChild(container);
  return res;
}
"""

# ---------------------------------------------------------------------------
# 主流程：逐頁擷取 + 組 PPTX（全程記憶體處理，不落地中繼檔）
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = px(PX_W)
prs.slide_height = px(PX_H)
blank = prs.slide_layouts[6]

with sync_playwright() as p:
    browser = p.chromium.launch(channel="chrome")
    page = browser.new_page(viewport={"width": 1200, "height": 1500}, device_scale_factor=1)

    for f in files:
        name = os.path.splitext(os.path.basename(f))[0]
        url = "file://" + os.path.abspath(f)
        page.goto(url, wait_until="networkidle")
        page.wait_for_timeout(200)  # 字型/排版安定

        units = page.evaluate(EXTRACT_JS)

        # 滿版頁規則：非封面、第一個物件貼著頁頂邊距開始的頁，首物件下移到 80pt
        page_class = page.evaluate("() => (document.querySelector('.page').className || '')")
        if units and "cover" not in page_class.split():
            first = min(units, key=lambda u: u["y"])
            first_y_target = FULLBLEED_FIRST_Y_PT / 0.75   # 80pt = 106.7px
            if first["y"] < first_y_target:
                first["y"] = first_y_target

        # .body 覆寫單元：字級/框寬改變前，先重新量測（螢光筆位置＋精確高度）
        for u in units:
            classes = (u.get("className") or "").split()
            if "body" in classes:
                remeasured = page.evaluate(REMEASURE_JS, {
                    "html": u["outerHTML"], "boxWpx": BODY_BOX_W_PX,
                    "fontPx": BODY_FONT_PX, "lineHeight": BODY_LINE_SPACING,
                    "letterSpacingEm": BODY_LETTER_SPACING,
                })
                u["_body_height"] = remeasured["height"]
                u["_body_hlRects"] = remeasured["hlRects"]  # 相對於新框左上角（局部座標）
            elif u.get("hasHlLine"):
                # 句級螢光筆：量套上新字距/行距後黃框的自然尺寸
                u["_hlline_size"] = page.evaluate(HLLINE_REMEASURE_JS, {
                    "html": u["outerHTML"],
                    "letterSpacingEm": HLLINE_LETTER_SPACING,
                    "lineHeight": HLLINE_LINE_SPACING,
                })

        # SVG 圖表獨立擷取為透明 PNG，之後從背景中隱藏
        svg_infos = page.evaluate(SVG_PREP_JS)
        svg_shots = []
        for i, info in enumerate(svg_infos):
            svg_el = page.query_selector(f'[data-export-svg="{i}"]')
            if svg_el is None:
                continue
            svg_shots.append((info, svg_el.screenshot(omit_background=True)))
        if svg_infos:
            page.evaluate(SVG_RESTORE_HIDE_JS)
            print(f"{name}: {len(svg_shots)} 個 SVG 圖表抽出為獨立物件")

        el = page.query_selector(".page")
        if el is None:
            print("略過（找不到 .page）:", f)
            continue
        bg_bytes = el.screenshot()  # 不給 path，直接拿 bytes，不落地
        bg_img = Image.open(io.BytesIO(bg_bytes))

        slide = prs.slides.add_slide(blank)
        top_color, dev_px = dominant_color_deviation(bg_img)
        if dev_px <= 50:
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*top_color)
            print(f"{name}: 背景判定為純色 {top_color}（偏差像素 {dev_px}），改用底色，不嵌圖")
        else:
            slide.shapes.add_picture(io.BytesIO(bg_bytes), 0, 0, width=px(PX_W), height=px(PX_H))

        # SVG 圖表：緊接在背景之後加入（z-order 在文字物件之下），各自是可移動的獨立圖片
        for info, shot in svg_shots:
            slide.shapes.add_picture(io.BytesIO(shot), px(info["x"]), px(info["y"]),
                                     width=px(info["w"]), height=px(info["h"]))

        for u in units:
            classes = (u.get("className") or "").split()
            is_body = "body" in classes
            is_slot = ("illust-slot" in classes) or ("photo-slot" in classes)
            is_hlline = bool(u.get("hasHlLine"))

            # 找出這個文字單元裡，非透明、非螢光黃的背景色（視為卡片/藥丸底色）
            block_bg = None
            for para in u["paragraphs"]:
                for r in para:
                    if not is_transparent(r["bg"]) and not is_highlight(r["bg"]):
                        block_bg = r["bg"]
                        break
                if block_bg:
                    break

            box_x, box_y, box_w, box_h = u["x"], u["y"], u["w"], u["h"]
            if is_body:
                # .body 段落覆寫：窄化置中框，y 不動，高度用重新量測的精確值
                box_x = BODY_BOX_X_PX
                box_w = BODY_BOX_W_PX
                box_h = u.get("_body_height", u["h"])
            elif is_hlline:
                # 句級螢光筆：黃框自己當文字框。以原黃框中心為錨，
                # 套新字距/行距後的自然尺寸置中擺放
                old = None
                for para in u["paragraphs"]:
                    for r in para:
                        if r.get("hlRects"):
                            old = r["hlRects"][0]
                            break
                    if old:
                        break
                if old is None:
                    old = {"x": u["x"], "y": u["y"], "w": u["w"], "h": u["h"]}
                size = u.get("_hlline_size") or {"w": old["w"], "h": old["h"]}
                pad_v = HLLINE_PAD_V_PT / 0.75   # 上下各 8pt 內距
                box_w = size["w"]
                box_h = size["h"] + 2 * pad_v
                box_x = old["x"] + old["w"] / 2 - box_w / 2
                box_y = old["y"] + old["h"] / 2 - box_h / 2

            # 螢光筆色塊：先畫（在文字框之前加入 shape tree，z-order 才會在文字底下）
            # （句級 .hl-line 除外——它的黃框就是文字框本身，不另畫色塊）
            if is_body:
                for hr in u.get("_body_hlRects") or []:
                    add_highlight_shape(slide, box_x + hr["x"], box_y + hr["y"], hr["w"], hr["h"])
            elif not is_hlline:
                for para in u["paragraphs"]:
                    for r in para:
                        for hr in (r.get("hlRects") or []):
                            add_highlight_shape(slide, hr["x"], hr["y"], hr["w"], hr["h"])

            # 一律用 textbox 承載文字（含帶底色的藥丸/卡片/插畫位），確保單一物件
            box = slide.shapes.add_textbox(px(box_x), px(box_y), px(box_w), px(box_h))
            tf = box.text_frame

            if is_slot:
                # 插畫位：直角（保單一物件）＋灰虛線框＋無陰影；底色照原樣（通常是白）
                if block_bg:
                    box.fill.solid()
                    box.fill.fore_color.rgb = rgb_str_to_hex(block_bg)
                box.line.color.rgb = RGBColor(*SLOT_BORDER_GRAY)
                box.line.width = Pt(1.5)
                box.line.dash_style = MSO_LINE_DASH_STYLE.DASH
                box.shadow.inherit = False
            elif is_hlline:
                # 句級螢光筆：黃底圓角文字框，單一物件
                set_round_rect(box, HLLINE_RADIUS_PX / max(1, min(box_w, box_h)) * 100000)
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(*HIGHLIGHT_YELLOW)
                box.line.fill.background()
                box.shadow.inherit = False
            elif block_bg:
                is_pill = u["h"] < 80 and u["w"] / max(u["h"], 1) > 2
                set_round_rect(box, 50000 if is_pill else 12000)
                box.fill.solid()
                box.fill.fore_color.rgb = rgb_str_to_hex(block_bg)
                box.line.fill.background()
                box.shadow.inherit = False

            tf.word_wrap = False if is_hlline else (u.get("lines", 2) != 1)
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE if (is_slot or is_hlline) else MSO_ANCHOR.TOP
            try:
                tf.auto_size = MSO_AUTO_SIZE.NONE
            except Exception:
                pass

            first_para = True
            for para_runs in u["paragraphs"]:
                if not para_runs:
                    continue
                para = tf.paragraphs[0] if first_para else tf.add_paragraph()
                first_para = False
                para.alignment = PP_ALIGN.CENTER if is_hlline else align_of(u["textAlign"])
                if is_body:
                    para.line_spacing = BODY_LINE_SPACING
                elif is_hlline:
                    para.line_spacing = HLLINE_LINE_SPACING
                for rd in para_runs:
                    text = rd["text"].strip()
                    if not text:
                        continue
                    run = para.add_run()
                    run.text = text
                    if is_slot:
                        run.font.size = Pt(SLOT_FONT_PT)
                    elif is_body:
                        run.font.size = Pt(BODY_FONT_PT)
                    else:
                        run.font.size = Pt(round(rd["fontSize"] * 0.75, 1))
                    if is_body:
                        # 字距：OOXML 的 spc 屬性，單位 1/100 pt（28pt × 13% = 3.64pt = 364）
                        run._r.get_or_add_rPr().set('spc', str(int(round(BODY_FONT_PT * BODY_LETTER_SPACING * 100))))
                    elif is_hlline:
                        pt_size = round(rd["fontSize"] * 0.75, 1)
                        run._r.get_or_add_rPr().set('spc', str(int(round(pt_size * HLLINE_LETTER_SPACING * 100))))
                    run.font.bold = False
                    fname = jinxuan_font_name(rd["weight"])
                    run.font.name = fname
                    set_cjk_font(run, fname)
                    col = rgb_str_to_hex(rd["color"])
                    if col:
                        run.font.color.rgb = col
                    # 注意：螢光筆已改用圖形色塊處理（見上方 add_highlight_shape），
                    # 這裡不再對 run 設定 <a:highlight>，避免跟色塊重複疊色。

        print(f"{name}: {len(units)} 個文字單元已加入")

    browser.close()

prs.save(out_path)
print("完成:", out_path)
